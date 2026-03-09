"""
PPT Parser Module
- Extract text, images, and tables from PowerPoint files
- Convert images to base64 for Claude Vision API
"""

import io
import base64
from typing import List, Dict, Any, Optional, Union
from pptx import Presentation
from pptx.util import Inches
from pptx.enum.shapes import MSO_SHAPE_TYPE
from PIL import Image


class PPTParseError(Exception):
    """PPT 파싱 관련 커스텀 예외"""

    def __init__(self, message: str, slide_num: Optional[int] = None):
        self.message = message
        self.slide_num = slide_num
        super().__init__(self.message)


def extract_slide_content(
    uploaded_file,
    include_notes: bool = True,
    skip_on_error: bool = True
) -> List[Dict[str, Any]]:
    """
    PPT 파일에서 모든 슬라이드 콘텐츠를 추출합니다.

    Args:
        uploaded_file: Streamlit UploadedFile 또는 파일 객체
        include_notes: 발표자 노트 포함 여부 (기본: True)
        skip_on_error: 에러 발생 시 해당 슬라이드 스킵 (기본: True)

    Returns:
        List[Dict[str, Any]]: 슬라이드별 추출된 콘텐츠

    Raises:
        PPTParseError: skip_on_error=False일 때 파싱 실패 시
    """
    try:
        prs = Presentation(uploaded_file)
    except Exception as e:
        raise PPTParseError(f"Failed to open PPT file: {str(e)}")

    slides_data = []

    for slide_num, slide in enumerate(prs.slides, start=1):
        slide_content = {
            "slide_num": slide_num,
            "texts": [],
            "images": [],
            "tables": [],
            "has_chart": False,
            "has_diagram": False,
            "notes": None,  # 새로 추가 (Phase 2에서 구현)
            "parse_errors": []  # 새로 추가: 파싱 중 발생한 에러 목록
        }

        try:
            for shape in slide.shapes:
                try:
                    # Extract text
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_content["texts"].append(shape.text.strip())

                    # Extract images
                    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                        image_data = extract_image(shape)
                        if image_data:
                            slide_content["images"].append(image_data)

                    # Extract tables
                    if shape.has_table:
                        table_data = extract_table(shape.table)
                        slide_content["tables"].append(table_data)

                    # Check for charts
                    if shape.has_chart:
                        slide_content["has_chart"] = True

                    # Check for grouped shapes (diagrams)
                    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                        slide_content["has_diagram"] = True

                except Exception as shape_error:
                    error_msg = f"Error processing shape in slide {slide_num}: {str(shape_error)}"
                    slide_content["parse_errors"].append(error_msg)
                    if not skip_on_error:
                        raise PPTParseError(error_msg, slide_num)

            # Extract presenter notes if requested
            if include_notes:
                try:
                    notes = extract_notes(slide)
                    if notes:
                        slide_content["notes"] = notes
                except Exception as notes_error:
                    error_msg = f"Error extracting notes from slide {slide_num}: {str(notes_error)}"
                    slide_content["parse_errors"].append(error_msg)
                    # 노트 추출 실패는 치명적이지 않으므로 skip_on_error와 관계없이 계속 진행

            slides_data.append(slide_content)

        except PPTParseError:
            raise  # Re-raise PPTParseError if skip_on_error=False
        except Exception as slide_error:
            error_msg = f"Error processing slide {slide_num}: {str(slide_error)}"
            if skip_on_error:
                # 에러가 발생한 슬라이드도 추가하되, 에러 정보를 포함
                slide_content["parse_errors"].append(error_msg)
                slides_data.append(slide_content)
            else:
                raise PPTParseError(error_msg, slide_num)

    return slides_data


def extract_notes(slide) -> Optional[str]:
    """
    슬라이드의 발표자 노트를 추출합니다.

    Args:
        slide: python-pptx Slide 객체

    Returns:
        Optional[str]: 발표자 노트 텍스트 또는 None

    Example:
        >>> notes = extract_notes(slide)
        >>> if notes:
        ...     print(f"Notes: {notes[:100]}...")
    """
    try:
        notes_slide = slide.notes_slide
        if notes_slide and notes_slide.notes_text_frame:
            notes_text = notes_slide.notes_text_frame.text.strip()
            return notes_text if notes_text else None
    except Exception:
        return None
    return None


def extract_image(shape) -> Optional[Dict[str, Any]]:
    """
    Extract image from shape and convert to base64.

    Args:
        shape: python-pptx Shape 객체

    Returns:
        Optional[Dict[str, Any]]: 이미지 데이터 또는 None (추출 실패 시)
    """
    try:
        image = shape.image
        image_bytes = image.blob
        image_ext = image.ext

        # Convert to base64
        base64_image = base64.b64encode(image_bytes).decode('utf-8')

        # Get image dimensions if possible
        width = shape.width.inches if hasattr(shape, 'width') else None
        height = shape.height.inches if hasattr(shape, 'height') else None

        return {
            "base64": base64_image,
            "format": image_ext,
            "width": width,
            "height": height,
            "media_type": get_media_type(image_ext)
        }
    except Exception:
        # 에러 발생 시 None 반환 (호출부에서 처리)
        return None


def extract_table(table) -> List[List[str]]:
    """
    테이블 콘텐츠를 2차원 배열로 추출합니다.

    Args:
        table: python-pptx Table 객체

    Returns:
        List[List[str]]: 2차원 문자열 배열
    """
    table_data = []
    try:
        for row in table.rows:
            row_data = []
            for cell in row.cells:
                row_data.append(cell.text.strip())
            table_data.append(row_data)
    except Exception:
        # 테이블 추출 실패 시 빈 테이블 반환
        return []
    return table_data


def get_media_type(ext: str) -> str:
    """
    Get MIME type from file extension.
    """
    media_types = {
        "png": "image/png",
        "jpg": "image/jpeg",
        "jpeg": "image/jpeg",
        "gif": "image/gif",
        "webp": "image/webp",
        "bmp": "image/bmp"
    }
    return media_types.get(ext.lower(), "image/png")


def get_slide_text_combined(
    slide_data: Dict[str, Any],
    include_notes: bool = False
) -> str:
    """
    슬라이드의 모든 텍스트를 하나의 문자열로 결합합니다.

    Args:
        slide_data: SlideContent 딕셔너리
        include_notes: 발표자 노트 포함 여부 (기본: False)

    Returns:
        str: 결합된 텍스트
    """
    texts = slide_data.get("texts", [])
    tables = slide_data.get("tables", [])
    notes = slide_data.get("notes")

    combined = "\n".join(texts)

    # Add table content
    for table in tables:
        table_text = "\n".join([" | ".join(row) for row in table])
        combined += f"\n[표]\n{table_text}"

    # Add presenter notes if requested and available
    if include_notes and notes:
        combined += f"\n\n[발표자 노트]\n{notes}"

    return combined


def get_all_text_content(slides_data: List[Dict[str, Any]]) -> str:
    """
    Get all text content from all slides as a single string.
    """
    all_text = []
    for slide in slides_data:
        slide_text = get_slide_text_combined(slide)
        if slide_text.strip():
            all_text.append(f"[슬라이드 {slide['slide_num']}]\n{slide_text}")

    return "\n\n".join(all_text)


def get_slides_with_images(slides_data: List[Dict[str, Any]]) -> List[Dict[str, Any]]:
    """
    Filter slides that contain images for Vision API processing.
    """
    return [slide for slide in slides_data if slide.get("images")]


def prepare_vision_content(slide_data: Dict[str, Any]) -> List[Dict[str, Any]]:
    """
    Prepare content list for Claude Vision API.
    Returns list of content blocks (text and images).
    """
    content = []

    # Add text context
    text_content = get_slide_text_combined(slide_data)
    if text_content:
        content.append({
            "type": "text",
            "text": f"슬라이드 {slide_data['slide_num']}의 텍스트 내용:\n{text_content}"
        })

    # Add images
    for img in slide_data.get("images", []):
        content.append({
            "type": "image",
            "source": {
                "type": "base64",
                "media_type": img["media_type"],
                "data": img["base64"]
            }
        })

    return content
