import streamlit as st
from pptx import Presentation
import io

# Page configuration
st.set_page_config(
    page_title="Smart Study Assistant",
    page_icon="📚",
    layout="wide",
    initial_sidebar_state="expanded"
)

# Initialize session state
if "processed" not in st.session_state:
    st.session_state.processed = False
if "slides_data" not in st.session_state:
    st.session_state.slides_data = []
if "summary" not in st.session_state:
    st.session_state.summary = None
if "quizzes" not in st.session_state:
    st.session_state.quizzes = []
if "quiz_answers" not in st.session_state:
    st.session_state.quiz_answers = {}
if "quiz_submitted" not in st.session_state:
    st.session_state.quiz_submitted = False
if "wrong_answers" not in st.session_state:
    st.session_state.wrong_answers = []
if "chat_history" not in st.session_state:
    st.session_state.chat_history = []
if "current_quiz_stage" not in st.session_state:
    st.session_state.current_quiz_stage = 0

# ============================================
# SIDEBAR: File Upload & Settings
# ============================================
with st.sidebar:
    st.header("📚 Smart Study Assistant")
    st.markdown("---")

    # File Uploader
    st.subheader("📁 파일 업로드")
    uploaded_file = st.file_uploader(
        "PPT 파일을 드래그하거나 클릭하여 업로드",
        type=["pptx"],
        help="PowerPoint 파일(.pptx)만 지원됩니다"
    )

    st.markdown("---")

    # Level Selector
    st.subheader("🎯 학습 설정")
    level = st.selectbox(
        "난이도 선택",
        options=["중학생", "고등학생", "대학생", "전문가"],
        index=2,
        help="선택한 수준에 맞춰 요약과 퀴즈가 생성됩니다"
    )

    # Quiz Configuration
    st.subheader("📝 퀴즈 설정")
    num_questions = st.slider(
        "문제 수",
        min_value=5,
        max_value=30,
        value=10,
        step=5
    )

    col1, col2 = st.columns(2)
    with col1:
        include_multiple_choice = st.checkbox("객관식", value=True)
    with col2:
        include_short_answer = st.checkbox("단답형", value=True)

    col3, col4 = st.columns(2)
    with col3:
        include_fill_blank = st.checkbox("빈칸 채우기", value=True)
    with col4:
        include_essay = st.checkbox("서술형", value=False)

    st.markdown("---")

    # Process Button
    process_btn = st.button(
        "🚀 학습 자료 생성 시작",
        type="primary",
        use_container_width=True,
        disabled=uploaded_file is None
    )

    if process_btn and uploaded_file:
        with st.spinner("AI가 학습 자료를 분석 중입니다..."):
            # TODO: Implement actual processing logic
            prs = Presentation(uploaded_file)
            slides_data = []

            for i, slide in enumerate(prs.slides):
                slide_content = {
                    "slide_num": i + 1,
                    "texts": [],
                    "images": [],
                    "vision_analysis": None
                }

                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_content["texts"].append(shape.text.strip())
                    # TODO: Extract images

                slides_data.append(slide_content)

            st.session_state.slides_data = slides_data
            st.session_state.processed = True

            # TODO: Generate summary and quizzes using Claude API
            st.session_state.summary = {
                "one_line": "이 자료는 [주제]에 대한 핵심 개념을 다룹니다.",
                "keywords": ["키워드1", "키워드2", "키워드3", "키워드4", "키워드5"]
            }

            # Placeholder quizzes
            st.session_state.quizzes = [
                {
                    "stage": "어휘다지기",
                    "questions": [
                        {
                            "id": 1,
                            "type": "multiple_choice",
                            "question": "[샘플] 다음 중 올바른 설명은?",
                            "options": ["보기 1", "보기 2", "보기 3", "보기 4"],
                            "answer": 0,
                            "source_slide": 1,
                            "explanation": "보기 1이 정답입니다."
                        }
                    ]
                },
                {
                    "stage": "실력다지기",
                    "questions": [
                        {
                            "id": 2,
                            "type": "short_answer",
                            "question": "[샘플] 핵심 개념을 한 단어로 답하세요.",
                            "answer": "정답",
                            "source_slide": 3,
                            "explanation": "슬라이드 3에서 설명된 내용입니다."
                        }
                    ]
                },
                {
                    "stage": "심화학습",
                    "questions": [
                        {
                            "id": 3,
                            "type": "essay",
                            "question": "[샘플] 본 자료의 핵심 논점을 서술하세요.",
                            "answer": "모범 답안 예시",
                            "source_slide": 5,
                            "explanation": "여러 슬라이드의 내용을 종합해야 합니다."
                        }
                    ]
                }
            ]

        st.success("✅ 학습 자료 생성 완료!")

    # Learning Progress
    if st.session_state.processed:
        st.markdown("---")
        st.subheader("📊 학습 진척도")

        total_questions = sum(len(stage["questions"]) for stage in st.session_state.quizzes)
        answered = len(st.session_state.quiz_answers)
        progress = answered / total_questions if total_questions > 0 else 0

        st.progress(progress)
        st.caption(f"퀴즈 진행: {answered}/{total_questions} 문제")

        if st.session_state.wrong_answers:
            wrong_count = len(st.session_state.wrong_answers)
            st.metric("오답 수", wrong_count)

# ============================================
# MAIN PANEL: Tabs
# ============================================
st.title("📚 Smart Study Assistant")
st.markdown("AI 기반 PPT 학습 도우미 - 핵심 정리, 퀴즈, 오답 노트, AI 튜터")

if not st.session_state.processed:
    st.info("👈 사이드바에서 PPT 파일을 업로드하고 '학습 자료 생성 시작' 버튼을 클릭하세요.")

    # Show placeholder content
    col1, col2, col3, col4 = st.columns(4)
    with col1:
        st.markdown("### 📊 Dashboard")
        st.markdown("PPT 요약 및 핵심 분석")
    with col2:
        st.markdown("### ✍️ Quiz Zone")
        st.markdown("맞춤형 퀴즈 풀이")
    with col3:
        st.markdown("### 📝 Review Note")
        st.markdown("오답 노트 & 피드백")
    with col4:
        st.markdown("### 🤖 AI Tutor")
        st.markdown("실시간 Q&A 챗봇")

else:
    # Main Tabs
    tab1, tab2, tab3, tab4 = st.tabs([
        "📊 Dashboard",
        "✍️ Quiz Zone",
        "📝 Review Note",
        "🤖 AI Tutor"
    ])

    # ============================================
    # TAB 1: Dashboard (학습자료 요약)
    # ============================================
    with tab1:
        st.header("📊 학습자료 요약")

        # Overall Summary Card
        st.subheader("📌 전체 요약")
        summary_col1, summary_col2 = st.columns([2, 1])

        with summary_col1:
            st.info(st.session_state.summary["one_line"])

        with summary_col2:
            st.markdown("**핵심 키워드**")
            keywords_html = " ".join([f"`{kw}`" for kw in st.session_state.summary["keywords"]])
            st.markdown(keywords_html)

        st.markdown("---")

        # Slide-by-slide Cards
        st.subheader("📑 슬라이드별 분석")

        for slide in st.session_state.slides_data:
            with st.expander(f"📄 Slide #{slide['slide_num']}", expanded=False):
                col1, col2 = st.columns([1, 1])

                with col1:
                    st.markdown("**📝 핵심 내용**")
                    if slide["texts"]:
                        for text in slide["texts"][:5]:  # Limit to 5 items
                            st.markdown(f"- {text[:200]}{'...' if len(text) > 200 else ''}")
                    else:
                        st.caption("텍스트 내용이 없습니다.")

                with col2:
                    st.markdown("**🖼️ Vision AI 분석**")
                    if slide["vision_analysis"]:
                        st.markdown(slide["vision_analysis"])
                    else:
                        st.caption("이미지 분석 결과가 없습니다.")

                    # Source button
                    if st.button(f"원본 슬라이드 보기", key=f"view_slide_{slide['slide_num']}"):
                        st.info(f"Slide #{slide['slide_num']} 원본 이미지 (구현 예정)")

    # ============================================
    # TAB 2: Quiz Zone (학습 확인)
    # ============================================
    with tab2:
        st.header("✍️ Quiz Zone")

        # Stage Progress
        stages = ["어휘다지기", "실력다지기", "심화학습"]
        current_stage = st.session_state.current_quiz_stage

        # Progress indicator
        progress_cols = st.columns(3)
        for i, stage in enumerate(stages):
            with progress_cols[i]:
                if i < current_stage:
                    st.success(f"✅ {stage}")
                elif i == current_stage:
                    st.info(f"▶️ {stage}")
                else:
                    st.markdown(f"⬜ {stage}")

        st.markdown("---")

        # Current Stage Questions
        if current_stage < len(st.session_state.quizzes):
            stage_data = st.session_state.quizzes[current_stage]
            st.subheader(f"📝 {stage_data['stage']}")

            # Progress Bar
            total_q = len(stage_data["questions"])
            answered_q = sum(1 for q in stage_data["questions"] if q["id"] in st.session_state.quiz_answers)
            st.progress(answered_q / total_q if total_q > 0 else 0)
            st.caption(f"진행: {answered_q}/{total_q} 문제")

            # Questions
            for q in stage_data["questions"]:
                st.markdown(f"**문제 {q['id']}** (출처: Slide #{q['source_slide']})")
                st.markdown(q["question"])

                q_key = f"q_{q['id']}"

                if q["type"] == "multiple_choice":
                    # Multiple choice buttons
                    selected = st.session_state.quiz_answers.get(q["id"])

                    cols = st.columns(len(q["options"]))
                    for i, option in enumerate(q["options"]):
                        with cols[i]:
                            btn_type = "primary" if selected == i else "secondary"

                            # Show result if answered
                            if selected is not None:
                                if i == q["answer"]:
                                    st.success(f"✅ {option}")
                                elif selected == i:
                                    st.error(f"❌ {option}")
                                else:
                                    st.button(option, key=f"{q_key}_opt_{i}", disabled=True)
                            else:
                                if st.button(option, key=f"{q_key}_opt_{i}"):
                                    st.session_state.quiz_answers[q["id"]] = i
                                    if i != q["answer"]:
                                        st.session_state.wrong_answers.append({
                                            "question": q,
                                            "user_answer": option,
                                            "correct_answer": q["options"][q["answer"]]
                                        })
                                    st.rerun()

                elif q["type"] == "short_answer":
                    user_answer = st.text_input("답변 입력", key=q_key)

                    if st.button("제출", key=f"{q_key}_submit"):
                        st.session_state.quiz_answers[q["id"]] = user_answer
                        if user_answer.strip().lower() != q["answer"].strip().lower():
                            st.session_state.wrong_answers.append({
                                "question": q,
                                "user_answer": user_answer,
                                "correct_answer": q["answer"]
                            })
                        st.rerun()

                    if q["id"] in st.session_state.quiz_answers:
                        if st.session_state.quiz_answers[q["id"]].strip().lower() == q["answer"].strip().lower():
                            st.success(f"✅ 정답입니다!")
                        else:
                            st.error(f"❌ 오답입니다. 정답: {q['answer']}")

                elif q["type"] == "essay":
                    user_answer = st.text_area("답변 작성", key=q_key, height=150)

                    if st.button("제출", key=f"{q_key}_submit"):
                        st.session_state.quiz_answers[q["id"]] = user_answer
                        # Essay questions need AI evaluation
                        st.info("서술형 답변이 제출되었습니다. AI 평가 기능은 구현 예정입니다.")

                st.markdown("---")

            # Stage Navigation
            col1, col2 = st.columns(2)
            with col1:
                if current_stage > 0:
                    if st.button("⬅️ 이전 단계"):
                        st.session_state.current_quiz_stage -= 1
                        st.rerun()
            with col2:
                if current_stage < len(stages) - 1:
                    if st.button("다음 단계 ➡️"):
                        st.session_state.current_quiz_stage += 1
                        st.rerun()
        else:
            st.success("🎉 모든 퀴즈를 완료했습니다!")
            if st.button("처음부터 다시 풀기"):
                st.session_state.current_quiz_stage = 0
                st.session_state.quiz_answers = {}
                st.session_state.wrong_answers = []
                st.rerun()

    # ============================================
    # TAB 3: Review Note (오답 노트 & 피드백)
    # ============================================
    with tab3:
        st.header("📝 Review Note")

        if not st.session_state.wrong_answers:
            st.success("🎉 오답이 없습니다! 훌륭합니다!")
        else:
            # Wrong Answers Table
            st.subheader("❌ 오답 목록")

            for i, wrong in enumerate(st.session_state.wrong_answers):
                with st.expander(f"문제 {wrong['question']['id']}: {wrong['question']['question'][:50]}...", expanded=True):
                    col1, col2, col3 = st.columns(3)

                    with col1:
                        st.markdown("**내 답변**")
                        st.error(wrong["user_answer"])

                    with col2:
                        st.markdown("**정답**")
                        st.success(wrong["correct_answer"])

                    with col3:
                        st.markdown("**출처**")
                        slide_num = wrong["question"]["source_slide"]
                        if st.button(f"📄 Slide #{slide_num}로 이동", key=f"goto_slide_{i}"):
                            st.info(f"Dashboard 탭의 Slide #{slide_num}을 확인하세요.")

                    # AI Explanation
                    st.markdown("**💡 AI 해설**")
                    st.info(wrong["question"]["explanation"])

            st.markdown("---")

            # Weakness Analysis (Placeholder)
            st.subheader("📊 취약점 분석")
            st.caption("학습 데이터가 충분히 쌓이면 레이더 차트가 표시됩니다.")

            # Placeholder for Plotly radar chart
            col1, col2 = st.columns([2, 1])
            with col1:
                st.markdown("""
                **분석 항목 (구현 예정)**
                - 어휘력
                - 개념 이해도
                - 수치 해석
                - 논리적 추론
                - 종합적 사고
                """)
            with col2:
                st.metric("총 오답 수", len(st.session_state.wrong_answers))
                total_q = sum(len(stage["questions"]) for stage in st.session_state.quizzes)
                accuracy = ((total_q - len(st.session_state.wrong_answers)) / total_q * 100) if total_q > 0 else 0
                st.metric("정답률", f"{accuracy:.1f}%")

    # ============================================
    # TAB 4: AI Tutor (실시간 Q&A)
    # ============================================
    with tab4:
        st.header("🤖 AI Tutor")
        st.markdown("PPT 내용에 대해 궁금한 점을 자유롭게 질문하세요.")

        # Suggested Questions
        st.subheader("💡 추천 질문")
        suggested_cols = st.columns(3)
        suggested_questions = [
            "이 주제의 핵심 개념이 뭐야?",
            "시험에 나올만한 중요 문장 3개만 뽑아줘",
            "이 내용을 쉽게 설명해줘"
        ]

        for i, sq in enumerate(suggested_questions):
            with suggested_cols[i]:
                if st.button(sq, key=f"suggested_{i}", use_container_width=True):
                    st.session_state.chat_history.append({"role": "user", "content": sq})
                    # TODO: Get AI response
                    st.session_state.chat_history.append({
                        "role": "assistant",
                        "content": f"'{sq}'에 대한 답변입니다. (AI 응답 기능 구현 예정)"
                    })
                    st.rerun()

        st.markdown("---")

        # Chat History
        chat_container = st.container()
        with chat_container:
            for msg in st.session_state.chat_history:
                if msg["role"] == "user":
                    st.chat_message("user").write(msg["content"])
                else:
                    st.chat_message("assistant").write(msg["content"])

        # Chat Input
        user_input = st.chat_input("질문을 입력하세요...")

        if user_input:
            st.session_state.chat_history.append({"role": "user", "content": user_input})

            # TODO: Implement actual AI response using LangChain
            ai_response = f"'{user_input}'에 대한 AI 튜터의 답변입니다. (실제 AI 응답 기능 구현 예정)"
            st.session_state.chat_history.append({"role": "assistant", "content": ai_response})
            st.rerun()

        # Clear Chat Button
        if st.session_state.chat_history:
            if st.button("🗑️ 대화 내역 지우기"):
                st.session_state.chat_history = []
                st.rerun()
